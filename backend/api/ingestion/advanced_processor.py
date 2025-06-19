"""
Advanced document processor with streaming, OCR, and batch processing capabilities.
Handles large documents, multiple formats, and provides real-time progress updates.
"""

import os
import sys
import io
import asyncio
import hashlib
import tempfile
import traceback
from typing import Dict, List, Tuple, Optional, Any, AsyncGenerator
from datetime import datetime
from pathlib import Path
import logging

import fitz  # PyMuPDF
import pytesseract
from PIL import Image
import numpy as np
import pandas as pd
from docx import Document as DocxDocument
from pptx import Presentation
try:
    import magic
except ImportError:
    magic = None
from channels.generic.websocket import AsyncJsonWebsocketConsumer
from django.conf import settings
from django.core.cache import cache
from celery import current_task
from dataclasses import dataclass, field
import aiofiles
import aiohttp

from .chunking_utils import chunk_text, chunk_thesis_by_chapter
from .figure_extractor import FigureExtractor
from .embeddings_utils import add_document_chunk_to_weaviate
from ..models import Document, Figure


logger = logging.getLogger(__name__)


@dataclass
class ProcessingProgress:
    """Track document processing progress."""
    total_pages: int = 0
    processed_pages: int = 0
    total_chunks: int = 0
    processed_chunks: int = 0
    current_stage: str = "initializing"
    errors: List[str] = field(default_factory=list)
    warnings: List[str] = field(default_factory=list)
    start_time: datetime = field(default_factory=datetime.now)
    
    @property
    def percentage(self) -> float:
        if self.total_pages == 0:
            return 0.0
        return (self.processed_pages / self.total_pages) * 100
    
    @property
    def elapsed_time(self) -> float:
        return (datetime.now() - self.start_time).total_seconds()
    
    def to_dict(self) -> Dict[str, Any]:
        return {
            "total_pages": self.total_pages,
            "processed_pages": self.processed_pages,
            "total_chunks": self.total_chunks,
            "processed_chunks": self.processed_chunks,
            "current_stage": self.current_stage,
            "percentage": self.percentage,
            "elapsed_time": self.elapsed_time,
            "errors": self.errors,
            "warnings": self.warnings
        }


class AdvancedDocumentProcessor:
    """
    Advanced document processor with support for:
    - Large document streaming
    - OCR for scanned documents
    - Multiple file formats
    - Progress tracking
    - Error recovery
    - Batch processing
    """
    
    SUPPORTED_FORMATS = {
        'application/pdf': 'pdf',
        'application/vnd.openxmlformats-officedocument.wordprocessingml.document': 'docx',
        'application/vnd.openxmlformats-officedocument.presentationml.presentation': 'pptx',
        'text/plain': 'txt',
        'text/markdown': 'md',
        'application/vnd.ms-powerpoint': 'ppt',
        'application/msword': 'doc'
    }
    
    def __init__(self, progress_callback=None, websocket_channel=None):
        self.progress = ProcessingProgress()
        self.progress_callback = progress_callback
        self.websocket_channel = websocket_channel
        self.ocr_languages = ['eng', 'fra', 'deu', 'spa']  # OCR languages
        
    async def process_document(
        self,
        file_path: str,
        metadata: Dict[str, Any],
        enable_ocr: bool = True,
        extract_figures: bool = True,
        chunk_size: int = 400,
        chunk_overlap: int = 100
    ) -> Tuple[Document, List[str]]:
        """
        Process a document with advanced features.
        
        Args:
            file_path: Path to the document
            metadata: Document metadata (title, author, etc.)
            enable_ocr: Enable OCR for scanned pages
            extract_figures: Extract figures and tables
            chunk_size: Target chunk size in words
            chunk_overlap: Overlap between chunks
            
        Returns:
            Tuple of (Document object, List of chunk IDs)
        """
        try:
            # Validate file
            if not os.path.exists(file_path):
                raise FileNotFoundError(f"File not found: {file_path}")
            
            # Detect file type
            file_type = await self._detect_file_type(file_path)
            if file_type not in self.SUPPORTED_FORMATS.values():
                raise ValueError(f"Unsupported file type: {file_type}")
            
            # Update progress
            await self._update_progress("validating", "File validated successfully")
            
            # Extract text based on file type
            if file_type == 'pdf':
                text, figures = await self._process_pdf(
                    file_path, enable_ocr, extract_figures
                )
            elif file_type == 'docx':
                text, figures = await self._process_docx(file_path)
            elif file_type == 'pptx':
                text, figures = await self._process_pptx(file_path)
            else:
                text = await self._process_text_file(file_path)
                figures = []
            
            # Create document record
            document = await self._create_document_record(metadata, file_path)
            
            # Chunk text
            chunks = await self._chunk_document(
                text, metadata.get('doc_type', 'paper'),
                chunk_size, chunk_overlap
            )
            
            # Add chunks to vector store
            chunk_ids = await self._add_chunks_to_vectorstore(
                chunks, document, metadata
            )
            
            # Process figures if extracted
            if figures:
                await self._process_figures(figures, document)
            
            await self._update_progress("completed", "Document processed successfully")
            return document, chunk_ids
            
        except Exception as e:
            logger.error(f"Error processing document: {str(e)}")
            self.progress.errors.append(str(e))
            await self._update_progress("failed", f"Error: {str(e)}")
            raise
    
    async def _detect_file_type(self, file_path: str) -> str:
        """Detect file type using python-magic."""
        if magic is None:
            # Fallback to file extension if magic is not available
            ext = Path(file_path).suffix.lower()
            ext_to_type = {
                '.pdf': 'pdf',
                '.docx': 'docx',
                '.pptx': 'pptx',
                '.txt': 'text',
                '.csv': 'csv'
            }
            return ext_to_type.get(ext, 'unknown')
        
        mime = magic.Magic(mime=True)
        mime_type = mime.from_file(file_path)
        return self.SUPPORTED_FORMATS.get(mime_type, 'unknown')
    
    async def _process_pdf(
        self, 
        file_path: str,
        enable_ocr: bool,
        extract_figures: bool
    ) -> Tuple[str, List[Dict]]:
        """Process PDF with streaming and OCR support."""
        text_parts = []
        figures = []
        
        try:
            doc = fitz.open(file_path)
            self.progress.total_pages = len(doc)
            
            for page_idx, page in enumerate(doc):
                await self._update_progress(
                    "extracting_text",
                    f"Processing page {page_idx + 1}/{self.progress.total_pages}"
                )
                
                # Extract text
                page_text = page.get_text()
                
                # If no text and OCR enabled, perform OCR
                if not page_text.strip() and enable_ocr:
                    page_text = await self._ocr_page(page)
                
                text_parts.append(page_text)
                
                # Extract figures if enabled
                if extract_figures:
                    page_figures = await self._extract_page_figures(page, page_idx)
                    figures.extend(page_figures)
                
                self.progress.processed_pages += 1
                
                # Yield control to prevent blocking
                await asyncio.sleep(0)
            
            doc.close()
            
        except Exception as e:
            logger.error(f"Error processing PDF: {str(e)}")
            raise
        
        return '\n'.join(text_parts), figures
    
    async def _ocr_page(self, page: fitz.Page) -> str:
        """Perform OCR on a PDF page."""
        try:
            # Convert page to image
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
            img_data = pix.tobytes("png")
            img = Image.open(io.BytesIO(img_data))
            
            # Perform OCR
            text = pytesseract.image_to_string(
                img,
                lang='+'.join(self.ocr_languages)
            )
            
            return text
            
        except Exception as e:
            logger.warning(f"OCR failed for page: {str(e)}")
            self.progress.warnings.append(f"OCR failed for a page: {str(e)}")
            return ""
    
    async def _process_docx(self, file_path: str) -> Tuple[str, List[Dict]]:
        """Process Word document."""
        text_parts = []
        figures = []
        
        try:
            doc = DocxDocument(file_path)
            self.progress.total_pages = len(doc.paragraphs)
            
            for para in doc.paragraphs:
                text_parts.append(para.text)
            
            # Extract images
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    # Extract image data
                    image_data = rel.target_part.blob
                    figures.append({
                        "type": "image",
                        "data": image_data,
                        "caption": ""  # DOCX doesn't have direct caption support
                    })
            
        except Exception as e:
            logger.error(f"Error processing DOCX: {str(e)}")
            raise
        
        return '\n'.join(text_parts), figures
    
    async def _process_pptx(self, file_path: str) -> Tuple[str, List[Dict]]:
        """Process PowerPoint presentation."""
        text_parts = []
        figures = []
        
        try:
            prs = Presentation(file_path)
            self.progress.total_pages = len(prs.slides)
            
            for slide_idx, slide in enumerate(prs.slides):
                await self._update_progress(
                    "extracting_text",
                    f"Processing slide {slide_idx + 1}/{self.progress.total_pages}"
                )
                
                # Extract text from slide
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_parts.append(shape.text)
                    
                    # Extract images
                    if shape.shape_type == 13:  # Picture
                        image = shape.image
                        figures.append({
                            "type": "image",
                            "data": image.blob,
                            "caption": f"Slide {slide_idx + 1}"
                        })
                
                self.progress.processed_pages += 1
                await asyncio.sleep(0)
            
        except Exception as e:
            logger.error(f"Error processing PPTX: {str(e)}")
            raise
        
        return '\n'.join(text_parts), figures
    
    async def _process_text_file(self, file_path: str) -> str:
        """Process plain text file."""
        async with aiofiles.open(file_path, 'r', encoding='utf-8') as f:
            return await f.read()
    
    async def _chunk_document(
        self,
        text: str,
        doc_type: str,
        chunk_size: int,
        chunk_overlap: int
    ) -> List[str]:
        """Chunk document text based on type."""
        await self._update_progress("chunking", "Creating document chunks")
        
        if doc_type == 'thesis':
            chunks = chunk_thesis_by_chapter(text)
        else:
            chunks = chunk_text(text, chunk_size, chunk_overlap)
        
        self.progress.total_chunks = len(chunks)
        return chunks
    
    async def _create_document_record(
        self,
        metadata: Dict[str, Any],
        file_path: str
    ) -> Document:
        """Create document record in database."""
        document = Document.objects.create(
            title=metadata.get('title', Path(file_path).stem),
            author=metadata.get('author', ''),
            doc_type=metadata.get('doc_type', 'paper'),
            year=metadata.get('year')
        )
        return document
    
    async def _add_chunks_to_vectorstore(
        self,
        chunks: List[str],
        document: Document,
        metadata: Dict[str, Any]
    ) -> List[str]:
        """Add chunks to vector store with progress tracking."""
        chunk_ids = []
        
        for idx, chunk in enumerate(chunks):
            await self._update_progress(
                "indexing",
                f"Indexing chunk {idx + 1}/{len(chunks)}"
            )
            
            chunk_metadata = {
                **metadata,
                "document_id": document.id,
                "chunk_index": idx,
                "total_chunks": len(chunks)
            }
            
            chunk_id = add_document_chunk_to_weaviate(chunk, chunk_metadata)
            if chunk_id:
                chunk_ids.append(chunk_id)
            
            self.progress.processed_chunks += 1
            await asyncio.sleep(0)
        
        return chunk_ids
    
    async def _extract_page_figures(
        self,
        page: fitz.Page,
        page_idx: int
    ) -> List[Dict]:
        """Extract figures from a PDF page."""
        figures = []
        
        # Get images
        image_list = page.get_images(full=True)
        
        for img_idx, img_info in enumerate(image_list):
            xref = img_info[0]
            
            try:
                # Extract image
                base_img = page.parent.extract_image(xref)
                
                figures.append({
                    "type": "image",
                    "page": page_idx + 1,
                    "data": base_img["image"],
                    "format": base_img["ext"],
                    "caption": ""
                })
                
            except Exception as e:
                logger.warning(f"Failed to extract image: {str(e)}")
        
        return figures
    
    async def _process_figures(
        self,
        figures: List[Dict],
        document: Document
    ) -> None:
        """Process and store extracted figures."""
        for idx, figure in enumerate(figures):
            await self._update_progress(
                "processing_figures",
                f"Processing figure {idx + 1}/{len(figures)}"
            )
            
            # Create figure record
            figure_obj = Figure.objects.create(
                document=document,
                figure_id=f"{document.id}_fig_{idx}",
                figure_type=figure.get('type', 'image'),
                caption=figure.get('caption', ''),
                page_number=figure.get('page', 1),
                metadata=figure
            )
            
            # Save figure file
            if 'data' in figure:
                # Save image data
                pass  # Implementation depends on storage backend
    
    async def _update_progress(
        self,
        stage: str,
        message: str = ""
    ) -> None:
        """Update processing progress."""
        self.progress.current_stage = stage
        
        # Call progress callback if provided
        if self.progress_callback:
            await self.progress_callback(self.progress.to_dict())
        
        # Send WebSocket update if channel provided
        if self.websocket_channel:
            await self.websocket_channel.send_json({
                "type": "progress",
                "data": self.progress.to_dict(),
                "message": message
            })
        
        # Update Celery task state if in Celery context
        if current_task:
            current_task.update_state(
                state='PROGRESS',
                meta=self.progress.to_dict()
            )


class BatchDocumentProcessor:
    """
    Process multiple documents in batch with error recovery.
    """
    
    def __init__(self, max_concurrent: int = 3):
        self.max_concurrent = max_concurrent
        self.results = []
        self.errors = []
    
    async def process_batch(
        self,
        file_paths: List[str],
        metadata_list: List[Dict[str, Any]],
        **kwargs
    ) -> Dict[str, Any]:
        """
        Process multiple documents concurrently.
        
        Args:
            file_paths: List of file paths
            metadata_list: List of metadata dicts
            **kwargs: Additional arguments for document processor
            
        Returns:
            Dict with results and errors
        """
        # Create semaphore for concurrency control
        semaphore = asyncio.Semaphore(self.max_concurrent)
        
        # Create tasks
        tasks = []
        for file_path, metadata in zip(file_paths, metadata_list):
            task = self._process_with_semaphore(
                semaphore, file_path, metadata, **kwargs
            )
            tasks.append(task)
        
        # Execute tasks
        await asyncio.gather(*tasks, return_exceptions=True)
        
        return {
            "total": len(file_paths),
            "successful": len(self.results),
            "failed": len(self.errors),
            "results": self.results,
            "errors": self.errors
        }
    
    async def _process_with_semaphore(
        self,
        semaphore: asyncio.Semaphore,
        file_path: str,
        metadata: Dict[str, Any],
        **kwargs
    ) -> None:
        """Process document with semaphore for concurrency control."""
        async with semaphore:
            try:
                processor = AdvancedDocumentProcessor()
                document, chunk_ids = await processor.process_document(
                    file_path, metadata, **kwargs
                )
                
                self.results.append({
                    "file_path": file_path,
                    "document_id": document.id,
                    "chunk_count": len(chunk_ids),
                    "status": "success"
                })
                
            except Exception as e:
                logger.error(f"Error processing {file_path}: {str(e)}")
                self.errors.append({
                    "file_path": file_path,
                    "error": str(e),
                    "traceback": traceback.format_exc(),
                    "status": "failed"
                })